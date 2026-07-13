from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent
DECK_PATH = OUT / "JOHN_HENRY_INVESTMENTS_PITCH_DECK.pptx"

NAVY = RGBColor(7, 16, 24)
GOLD = RGBColor(215, 173, 91)
GREEN = RGBColor(97, 211, 148)
WHITE = RGBColor(245, 248, 250)
MUTED = RGBColor(120, 138, 150)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.25), Inches(0.65))
    frame = title_box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(11.7), Inches(0.35))
        sub_frame = sub.text_frame
        sub_frame.text = subtitle
        sub_frame.paragraphs[0].runs[0].font.size = Pt(12)
        sub_frame.paragraphs[0].runs[0].font.color.rgb = MUTED


def add_footer(slide, number: int) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.08), Inches(12.0), Inches(0.25))
    frame = box.text_frame
    frame.text = f"John Henry Investments Platform | Confidential planning draft | {number}"
    p = frame.paragraphs[0]
    p.font.size = Pt(8)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def add_bullets(slide, bullets: list[str], x=0.75, y=1.35, w=5.8, h=5.2, size=16) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    for idx, text in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)


def add_metric(slide, label: str, value: str, x: float, y: float, w=2.35) -> None:
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(1.05))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(17, 33, 48)
    box.line.color.rgb = GOLD
    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = value
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = GREEN
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = WHITE


def add_table(slide, rows: list[list[str]], x, y, w, h, font_size=10) -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = GOLD if r_idx == 0 else RGBColor(17, 33, 48)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.bold = r_idx == 0
                p.font.color.rgb = NAVY if r_idx == 0 else WHITE


def add_bar_chart(slide, title: str, categories: list[str], values: list[float], x, y, w, h) -> None:
    data = CategoryChartData()
    data.categories = categories
    data.add_series(title, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data
    ).chart
    chart.has_legend = False
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.plots[0].data_labels.font.size = Pt(8)


def add_line_chart(slide, title: str, categories: list[str], values: list[float], x, y, w, h) -> None:
    data = CategoryChartData()
    data.categories = categories
    data.add_series(title, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), data
    ).chart
    chart.has_legend = False
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)


def add_slide(prs, number: int, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, title, subtitle)
    add_footer(slide, number)
    return slide


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    slide = add_slide(prs, 1, "John Henry Investments Platform", "AI-powered investment intelligence and family office operating system")
    add_metric(slide, "Base-case DCF value", "$133.5M", 0.8, 1.6)
    add_metric(slide, "Year 5 revenue", "$133.6M", 3.35, 1.6)
    add_metric(slide, "Year 5 EBITDA", "$53.2M", 5.9, 1.6)
    add_metric(slide, "Target platform users", "100K", 8.45, 1.6)
    add_bullets(slide, [
        "Subscription SaaS for investors, acquisition entrepreneurs, family offices, and professional advisors.",
        "Combines opportunity discovery, AI due diligence, macro intelligence, portfolio management, accounting, CRM, and reporting.",
        "Planning package includes pitch deck, DCF model, revenue/expenditure/marketing projections, personnel chart, and risk recommendations.",
    ], x=0.9, y=3.1, w=11.5, size=17)

    # 2
    slide = add_slide(prs, 2, "Problem", "Investment and acquisition workflows are fragmented")
    add_bullets(slide, [
        "Investors use disconnected tools for research, diligence, reporting, accounting, and CRM.",
        "Business buyers lack integrated SBA, valuation, document review, and debt-service analysis.",
        "Family offices and advisors need repeatable intelligence, governance, portfolio, and report workflows.",
        "AI outputs need structure, scoring, auditability, source discipline, and compliance controls.",
    ], w=11.5)

    # 3
    slide = add_slide(prs, 3, "Solution", "A unified investment intelligence platform")
    add_bullets(slide, [
        "John Henry Opportunity Score across stocks, businesses, crypto, real estate, and macro conditions.",
        "AI due diligence center for P&L, tax returns, balance sheets, bank statements, and fraud/risk flags.",
        "Business acquisition engine with EBITDA, SDE, DSCR, SBA qualification, and valuation models.",
        "Accounting, audit, CRM, banking/vendor, Microsoft Office, and reporting workflows in the backend.",
    ], w=11.5)

    # 4
    slide = add_slide(prs, 4, "Product modules", "Platform by module")
    modules = [
        ["Core", "Users, MFA, roles, dashboard"],
        ["Opportunity", "Discovery, acquisition engine, due diligence"],
        ["Intelligence", "Macro dashboard, weekly reports, AI assistant"],
        ["Wealth", "Portfolio tracking, wealth projection"],
        ["Business owner", "Governance, capital raising"],
        ["Operating system", "Accounting, CRM, integrations, reporting"],
    ]
    add_table(slide, [["Layer", "Modules"]] + modules, 0.8, 1.35, 11.8, 4.9, font_size=13)

    # 5
    slide = add_slide(prs, 5, "Revenue model", "B2C + B2B + enterprise subscription plans")
    add_table(slide, [
        ["Plan", "Price", "Primary customer"],
        ["Consumer", "$50/mo", "Retail investors and wealth builders"],
        ["Professional", "$299/mo", "Acquisition entrepreneurs and advisors"],
        ["Enterprise", "$1,500+/mo", "Family offices, firms, CPAs, attorneys, bankers"],
    ], 0.8, 1.4, 11.8, 2.4, font_size=13)
    add_bullets(slide, [
        "Expansion: annual contracts, branded reports, AI diligence jobs, integrations, API access, and white-label reporting.",
        "Enterprise ACH/wire billing can reduce card processing cost and improve net margin.",
    ], x=1, y=4.4, w=11, size=16)

    # 6 revenue chart
    slide = add_slide(prs, 6, "Revenue projection", "Base-case blended SaaS mix")
    add_bar_chart(slide, "Revenue ($M)", ["Y1", "Y2", "Y3", "Y4", "Y5"], [1.336, 6.682, 20.047, 66.822, 133.644], 0.8, 1.4, 11.7, 4.8)

    # 7 expenditure chart
    slide = add_slide(prs, 7, "Expenditure projection", "Platform operations, staffing, legal, marketing, payment processing")
    add_bar_chart(slide, "Opex ($M)", ["Y1", "Y2", "Y3", "Y4", "Y5"], [1.384, 3.710, 9.451, 34.605, 80.409], 0.8, 1.4, 11.7, 4.8)

    # 8 marketing chart
    slide = add_slide(prs, 8, "Marketing projection", "Spend and blended CAC")
    add_line_chart(slide, "Marketing ($M)", ["Y1", "Y2", "Y3", "Y4", "Y5"], [0.180, 0.750, 2.250, 8.000, 14.000], 0.8, 1.35, 6.0, 4.7)
    add_table(slide, [
        ["Year", "New users", "Blended CAC"],
        ["Y1", "1,000", "$180"],
        ["Y2", "4,000", "$188"],
        ["Y3", "10,000", "$225"],
        ["Y4", "35,000", "$229"],
        ["Y5", "50,000", "$280"],
    ], 7.15, 1.55, 5.2, 4.2, font_size=11)

    # 9 EBITDA
    slide = add_slide(prs, 9, "EBITDA projection", "Base-case after staffing, marketing, legal, platform ops, and payment processing")
    add_bar_chart(slide, "EBITDA ($M)", ["Y1", "Y2", "Y3", "Y4", "Y5"], [-0.048, 2.972, 10.595, 32.217, 53.235], 0.8, 1.4, 11.7, 4.8)

    # 10 DCF
    slide = add_slide(prs, 10, "Discounted cash flow", "Base-case planning valuation")
    add_metric(slide, "Discount rate", "18%", 0.8, 1.35)
    add_metric(slide, "Terminal growth", "3%", 3.35, 1.35)
    add_metric(slide, "Terminal value", "$233.7M", 5.9, 1.35)
    add_metric(slide, "DCF enterprise value", "$133.5M", 8.45, 1.35)
    add_bullets(slide, [
        "DCF uses FCF after cash taxes, capex, and working-capital needs.",
        "Model is intentionally more conservative than ARR-multiple valuation.",
        "Upside depends on enterprise pricing, lower CAC, lower AI cost, annual contracts, and retention.",
    ], x=0.9, y=3.05, w=11.5, size=17)

    # 11 org
    slide = add_slide(prs, 11, "Personnel hierarchy", "Sustainable organization design")
    add_metric(slide, "Year 1 team", "3-7", 0.7, 1.2)
    add_metric(slide, "Year 3 team", "20-45", 3.0, 1.2)
    add_metric(slide, "Year 4 team", "75-160", 5.3, 1.2)
    add_metric(slide, "Year 5 team", "150-350", 7.6, 1.2)
    add_bullets(slide, [
        "Founder / CEO",
        "Product + Design",
        "Engineering",
        "AI / Data / Research",
        "Security + DevOps",
        "Operations + Customer Success",
        "Sales + Partnerships",
        "Finance + Accounting",
        "Legal + Compliance",
    ], x=0.9, y=2.75, w=11, size=15)

    # 12 moat
    slide = add_slide(prs, 12, "Competitive moat", "Data, workflow, and scoring defensibility")
    add_bullets(slide, [
        "Proprietary John Henry Opportunity Score and explainability layer.",
        "AI due diligence engine learns from structured acquisition and financial workflows.",
        "Integrated accounting, CRM, portfolio, reports, and Office exports reduce switching.",
        "Enterprise data, report history, and diligence workspaces create compounding value.",
    ], w=11.5)

    # 13 GTM
    slide = add_slide(prs, 13, "Go-to-market", "Start with higher-value professional users")
    add_bullets(slide, [
        "Phase 1: founder-led controlled beta with acquisition entrepreneurs and advisors.",
        "Phase 2: professional plan launch with due diligence and acquisition workflows.",
        "Phase 3: CPA, attorney, banker, and family office referral partnerships.",
        "Phase 4: enterprise pilots with branded reports, integrations, and team accounts.",
        "Phase 5: consumer subscription scale after proof of retention and support economics.",
    ], w=11.5)

    # 14 risk
    slide = add_slide(prs, 14, "Risk controls", "Required before broad paid launch")
    add_bullets(slide, [
        "Financial, tax, legal, accounting, and AI-output disclaimers.",
        "Privacy policy, terms, data retention, and integration consent language.",
        "Role-based access control, audit logs, secure storage, and webhook verification.",
        "Human review for high-risk due diligence and investment-related outputs.",
    ], w=11.5)

    # 15 use funds
    slide = add_slide(prs, 15, "Use of funds", "Milestone-based capital allocation")
    add_table(slide, [
        ["Category", "Use"],
        ["Product engineering", "Auth, billing, database, AI, scoring, document workflows"],
        ["Infrastructure", "PostgreSQL, storage, workers, monitoring, security"],
        ["Compliance", "Legal review, privacy, disclaimers, advisory analysis"],
        ["Go-to-market", "Beta customers, professional launch, partnerships"],
        ["Operations", "Support, finance, reporting, account management"],
    ], 0.75, 1.3, 11.9, 4.9, font_size=12)

    # 16 ask
    slide = add_slide(prs, 16, "Investor ask", "Recommended structure to finalize")
    add_bullets(slide, [
        "Define amount to raise around milestones, not only valuation.",
        "Fund production MVP, legal/compliance package, Stripe, database persistence, AI scoring, document workflows, and pilot customers.",
        "Potential instruments: SAFE, convertible note, priced equity, or strategic investment.",
        "Next diligence items: live spreadsheet, customer discovery, legal memo, product screenshots, pilot LOIs.",
    ], w=11.5)

    # 17 extra recommendations
    slide = add_slide(prs, 17, "Additional recommendations", "Items investors will expect")
    add_bullets(slide, [
        "Churn, retention, and net revenue retention model.",
        "Customer acquisition channel model and sales pipeline.",
        "Competitive landscape and acquirer universe.",
        "Compliance roadmap and risk memo.",
        "Data moat and scoring methodology appendix.",
        "Product screenshots, demo video, and pilot customer quotes.",
    ], w=11.5)

    # 18 close
    slide = add_slide(prs, 18, "Closing", "From investment company to fintech operating system")
    add_bullets(slide, [
        "John Henry Investments Platform can become an AI-powered investment research, acquisition intelligence, due diligence, and family office operating system.",
        "The strongest assets are the proprietary score, diligence engine, report workflows, and integrated financial data layer.",
        "Immediate next step: define raise amount, pilot customers, compliance posture, and production MVP roadmap.",
    ], w=11.5)

    prs.save(DECK_PATH)
    print(f"Generated {DECK_PATH}")


if __name__ == "__main__":
    build_deck()
